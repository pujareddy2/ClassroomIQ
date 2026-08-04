from __future__ import annotations

from pathlib import Path

from pptx import Presentation  # type: ignore

from app.services.document_extractor.exceptions import CorruptedDocumentError


class PptExtractor:
    """Extract text from PPTX slides, titles, bullets, and notes."""

    def extract(self, file_path: str | Path) -> str:
        try:
            presentation = Presentation(str(file_path))
        except Exception as exc:  # pragma: no cover - dependency/runtime path
            raise CorruptedDocumentError("Unable to read the PPTX file") from exc

        parts: list[str] = []
        for slide in presentation.slides:
            slide_parts: list[str] = []

            # Title
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_parts.append(slide.shapes.title.text.strip())

            # Body text from shapes
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        slide_parts.append(text)

            # Speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                if notes_text_frame:
                    notes_text = notes_text_frame.text.strip()
                    if notes_text:
                        slide_parts.append(f"Notes: {notes_text}")

            if slide_parts:
                parts.append("\n".join(slide_parts))

        return "\n\n".join(parts)
