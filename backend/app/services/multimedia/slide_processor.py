"""
Slide & Presentation Deck Processor.
Extracts text, slide titles, notes, and visual page previews from PPTX and PDF lecture slides.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


class SlideProcessor:
    """Extracts text content and preview images from lecture presentation decks."""

    def process_presentation(
        self,
        presentation_path: Path,
        output_dir: Path,
    ) -> List[Dict[str, Any]]:
        """
        Parses a PPTX or PDF lecture presentation and extracts slide text + preview images.
        Returns a list of slide dicts: slide_number, title, text_content, preview_filename.
        """
        if not presentation_path.exists():
            raise FileNotFoundError(f"Presentation file not found: {presentation_path}")

        output_dir.mkdir(parents=True, exist_ok=True)
        ext = presentation_path.suffix.lower()

        if ext in {".pptx", ".ppt"}:
            return self._process_pptx(presentation_path, output_dir)
        elif ext == ".pdf":
            return self._process_pdf(presentation_path, output_dir)
        else:
            logger.warning("Unsupported presentation format %s for slide extraction", ext)
            return []

    def _process_pptx(self, pptx_path: Path, output_dir: Path) -> List[Dict[str, Any]]:
        """Extracts text, titles, and notes from a PowerPoint (.pptx) file."""
        slides_data: List[Dict[str, Any]] = []
        try:
            from pptx import Presentation
            prs = Presentation(str(pptx_path))

            for idx, slide in enumerate(prs.slides, start=1):
                title = ""
                text_pieces: List[str] = []

                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text and text != title:
                                text_pieces.append(text)

                # Extract speaker notes if any
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        text_pieces.append(f"[Notes: {notes}]")

                full_text = "\n".join(text_pieces)
                slides_data.append({
                    "slide_number": idx,
                    "title": title or f"Slide {idx}",
                    "text_content": full_text,
                    "preview_filename": None,  # PPTX preview requires libreoffice or pdf conversion
                })

            logger.info("Extracted %d slides from PPTX: %s", len(slides_data), pptx_path.name)
            return slides_data

        except ImportError:
            logger.warning("python-pptx not installed; skipping PPTX parsing.")
            return []
        except Exception as exc:
            logger.exception("Failed to parse PPTX file %s: %s", pptx_path, exc)
            return []

    def _process_pdf(self, pdf_path: Path, output_dir: Path) -> List[Dict[str, Any]]:
        """Extracts text and renders page images from a PDF presentation deck."""
        slides_data: List[Dict[str, Any]] = []

        # Try PyMuPDF (fitz) first for text and high-res rendering
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(pdf_path))

            for idx, page in enumerate(doc, start=1):
                page_text = page.get_text("text").strip()
                lines = [line.strip() for line in page_text.splitlines() if line.strip()]
                title = lines[0] if lines else f"Slide {idx}"
                body = "\n".join(lines[1:]) if len(lines) > 1 else ""

                # Render page image to PNG
                preview_filename = f"slide_{idx:03d}.png"
                preview_path = output_dir / preview_filename
                try:
                    pix = page.get_pixmap(dpi=150)
                    pix.save(str(preview_path))
                except Exception as render_err:
                    logger.debug("Failed to render page %d image: %s", idx, render_err)
                    preview_filename = None

                slides_data.append({
                    "slide_number": idx,
                    "title": title,
                    "text_content": body,
                    "preview_filename": preview_filename,
                })

            doc.close()
            logger.info("Extracted and rendered %d PDF slides from %s", len(slides_data), pdf_path.name)
            return slides_data

        except ImportError:
            pass
        except Exception as exc:
            logger.debug("PyMuPDF extraction failed: %s", exc)

        # Fallback to pdfplumber
        try:
            import pdfplumber
            with pdfplumber.open(str(pdf_path)) as pdf:
                for idx, page in enumerate(pdf.pages, start=1):
                    page_text = (page.extract_text() or "").strip()
                    lines = [line.strip() for line in page_text.splitlines() if line.strip()]
                    title = lines[0] if lines else f"Slide {idx}"
                    body = "\n".join(lines[1:]) if len(lines) > 1 else ""
                    slides_data.append({
                        "slide_number": idx,
                        "title": title,
                        "text_content": body,
                        "preview_filename": None,
                    })
            return slides_data
        except ImportError:
            pass
        except Exception as exc:
            logger.exception("Failed to parse PDF slides with pdfplumber: %s", exc)

        return slides_data
